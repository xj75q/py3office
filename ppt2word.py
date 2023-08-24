# coding=utf-8
from multiprocessing import Process
import sys
import importlib
import argparse
from pptx import Presentation
from docx import Document

importlib.reload(sys)
import re


class PPT2Word():
    def __init__(self, input_file):
        self.input = input_file
        self.save_name = r'new.docx'

    def run(self):
        wordfile = Document()
        try:
            # 给定ppt文件所在的路径
            pptx = Presentation(self.input)
            # 遍历ppt文件的所有幻灯片页
            for slide in pptx.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text_frame = shape.text_frame
                        for paragraph in text_frame.paragraphs:
                            content = paragraph.text
                            s = re.sub(u"[\\x00-\\x08\\x0b\\x0e-\\x1f\\x7f]", "", content)
                            wordfile.add_paragraph(s)

            wordfile.save(self.save_name)

        except Exception as err:
            print(err)


def parse_parm():
    parser = argparse.ArgumentParser()
    parser.add_argument('-i', '--input', help='input ppt file', required=True,
                        type=str)  # , default=False ,,action='store_true'
    args = parser.parse_args()
    input_str = args.input
    return input_str


if __name__ == "__main__":
    input_file = parse_parm()

    pw = PPT2Word(input_file)
    pw.run()
